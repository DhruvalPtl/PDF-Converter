import logging
import PyPDF2
import pdfkit
import win32com.client
from pdf2docx import Converter
from win32com import client
import os
import shutil
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import pdfplumber
import img2pdf
import fitz

class PDFConverter():
    def __init__(self):
        self.input_path_dir = r"C:/PdfConverter/Temp"
        # self.input_path_dir = os.path.join(sys._MEIPASS, "TEMP")        
        self.wkhtmltopdf_path = r"C:/PdfConverter/wkhtmltopdf/bin/wkhtmltopdf.exe"
        # self.wkhtmltopdf_path = os.path.join(sys._MEIPASS, "wkhtmltopdf/bin/wkhtmltopdf.exe")
        self.output_path = r"C:/PdfConverter/Temp1"
        # self.output_path = os.path.join(sys._MEIPASS, "TEMP1")
    
    def pdfmerger(self):
        pdf_files = [i for i in os.listdir(self.input_path_dir) if  i.lower().endswith(".pdf")]
        merger = PyPDF2.PdfMerger()

        for pdf in pdf_files:
            pdf_path = os.path.join(self.input_path_dir, pdf)
            with open(pdf_path, 'rb') as b:
                merger.append(b)
        merger.write(self.output_path + f"/{pdf}")

    def pdfcompress(self):
        pdf_files = [i for i in os.listdir(self.input_path_dir) if  i.lower().endswith(".pdf")]
        for pdf in pdf_files:
            input_pdf_path = os.path.join(self.input_path_dir, pdf)
            output_pdf_path = os.path.join(self.output_path, pdf)
            with open(input_pdf_path, 'rb') as b:
                pdf_reader = PyPDF2.PdfReader(b)
                pdf_writer = PyPDF2.PdfWriter()

                for page in pdf_reader.pages:
                    page.compress_content_streams()
                    pdf_writer.add_page(page)
            
                with open(output_pdf_path, 'wb') as output_pdf:
                    pdf_writer.write(output_pdf)

    def split_range(self, start_page, end_page):
        try:
            pdf_files = [i for i in os.listdir(self.input_path_dir) if  i.lower().endswith(".pdf")]
            for file_name in pdf_files:
                input_pdf_path = os.path.join(self.input_path_dir, file_name)
                
                if not os.path.exists(input_pdf_path):
                    logging.info(f"File not found: {input_pdf_path}")
                    continue

                with open(input_pdf_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    range_writer = PyPDF2.PdfWriter()

                    for page_num in range(start_page - 1, min(end_page, len(reader.pages))):
                        range_writer.add_page(reader.pages[page_num])

                    range_writer_output = os.path.join(self.output_path, f'{os.path.splitext(file_name)[0]}_range_{start_page}_{end_page}.pdf')
                    
                    with open(range_writer_output, 'wb') as output_pdf:
                        range_writer.write(output_pdf)
                    logging.info(f"Successfully saved range to {range_writer_output}")

        except Exception as e:
            logging.error(f"Error splitting PDF range: {e}")
            
    def extract_pages(self, page_numbers):
        try:
            pdf_files = [i for i in os.listdir(self.input_path_dir) if  i.lower().endswith(".pdf")]
            for file_name in pdf_files:
                input_pdf_path = os.path.join(self.input_path_dir, file_name)
                
                if not os.path.exists(input_pdf_path):
                    logging.info(f"File not found: {input_pdf_path}")
                    continue

                with open(input_pdf_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    extract_writer = PyPDF2.PdfWriter()

                    for page_num in page_numbers:
                        page_index = page_num - 1
                        if page_index < len(reader.pages):
                            extract_writer.add_page(reader.pages[page_index])
                        else:
                            logging.info(f"Page {page_num} is out of range")

                    extract_writer_output = os.path.join(self.output_path, f'{os.path.splitext(file_name)[0]}_extracted.pdf')
                    with open(extract_writer_output, 'wb') as output_pdf:
                        extract_writer.write(output_pdf)
                    logging.info(f"Successfully saved extracted pages to {extract_writer_output}")

        except Exception as e:
            logging.error(f"Error extracting PDF pages: {e}")
            
    def single_pages(self):
        pdf_files = [i for i in os.listdir(self.input_path_dir) if  i.lower().endswith(".pdf")]
        try:
            for file_name in pdf_files:
                input_pdf_path = os.path.join(self.input_path_dir, file_name)
                
                with open(input_pdf_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)

                    for page_num, page in enumerate(reader.pages):
                        single_page_writer = PyPDF2.PdfWriter()
                        single_page_writer.add_page(page)
                        single_page_output = os.path.join(self.output_path, f'{os.path.splitext(file_name)[0]}_page_{page_num + 1}.pdf')
                        
                        with open(single_page_output, 'wb') as output_pdf:
                            single_page_writer.write(output_pdf)
        except Exception as e:
            logging.info(f"Error splitting PDF into single pages: {e}")

    def emptydir(self):
        try:
            if os.path.exists(self.input_path_dir):
                for item in os.listdir(self.input_path_dir):
                    item_path = os.path.join(self.input_path_dir, item)
                    if os.path.isdir(item_path):
                        shutil.rmtree(item_path, ignore_errors=True)  
                    else:
                        os.remove(item_path)
                logging.info(f"Contents of directory {self.input_path_dir} have been removed.")
            else:
                logging.error(f"Directory {self.input_path_dir} does not exist.")
        except Exception as e:
            logging.error(f"Error removing contents of directory: {e}")

    def html2pdf(self):
        if not os.path.exists(self.output_path):
            os.makedirs(self.output_path)

        config = pdfkit.configuration(wkhtmltopdf=self.wkhtmltopdf_path)

        html_files = [i for i in os.listdir(self.input_path_dir) if i.lower().endswith(".html")]
        for i in html_files:
            input_html = os.path.join(self.input_path_dir, i)
            output_pdf = os.path.join(self.output_path, f"{os.path.splitext(i)[0]}.pdf")
            try:
                pdfkit.from_file(input_html, output_pdf, configuration=config)
                logging.info(f"Converted {i} to {output_pdf} successfully.")
            except Exception as e:
                logging.error(f"Failed to convert {i} to PDF: {e}")
    
    def pdf2word(self):
        try:
            pdf_files = [i for i in os.listdir(self.input_path_dir) if i.lower().endswith(".pdf")]
            if not pdf_files:
                logging.info("No PDF files found in the input directory.")
                return

            for pdf_file in pdf_files:
                pdf_path = os.path.join(self.input_path_dir, pdf_file)
                output_word = os.path.join(self.output_path, os.path.splitext(pdf_file)[0] + '.docx')
                cv = Converter(pdf_path)
                cv.convert(output_word, start=0, end=None)
                cv.close()
                logging.info(f"Converted {pdf_file} to Word successfully.")
        except Exception as e:
            logging.error(f"Error during PDF to Word conversion: {e}")
            
    def word2pdf(self):
        # Initialize the Word application
        word = win32com.client.Dispatch("Word.Application")
        
        try:
            # Ensure the output directory exists
            if not os.path.exists(self.output_path):
                os.makedirs(self.output_path)
            
            # List all .docx files in the input directory
            word_files = [f for f in os.listdir(self.input_path_dir) if f.lower().endswith('.docx')]
            
            if not word_files:
                logging.info("No .docx files found for conversion.")
                return

            for word_file in word_files:
                input_word = os.path.join(self.input_path_dir, word_file)
                output_pdf = os.path.join(self.output_path, os.path.splitext(word_file)[0] + '.pdf')
                
                input_word = os.path.normpath(input_word)
                output_pdf = os.path.normpath(output_pdf)
            
                logging.info(f"Converting {input_word} to {output_pdf}")
                try:
                    # Open the document
                    doc = word.Documents.Open(input_word)
                    
                    # Save the document as PDF
                    doc.SaveAs(output_pdf, FileFormat=17)  # 17 corresponds to PDF format
                    
                    # Close the document
                    doc.Close()
                    
                    logging.info(f"Conversion successful for {input_word}")
                except Exception as e:
                    logging.info(f"Error during conversion: {e}")
                    
        except Exception as e:
            logging.info(f"Error in word2pdf function: {e}")
        finally:
            # Quit the Word application
            word.Quit()

    def excel2pdf(self):
        try:
            excel_files = [f for f in os.listdir(self.input_path_dir) if f.lower().endswith(".xlsx")]
            
            if not excel_files:
                logging.info("No Excel files found in the input directory.")
                return

            config = pdfkit.configuration(wkhtmltopdf=self.wkhtmltopdf_path)
            if not config:
                logging.info(f"wkhtmltopdf file not found in {self.wkhtmltopdf_path}")
                return
            
            for file_name in excel_files:
                excel_file_path = os.path.join(self.input_path_dir, file_name)
                df = pd.read_excel(excel_file_path)
                html_file_path = os.path.join(self.input_path_dir, "Temp.html")
                df.to_html(html_file_path)
                output_pdf_path = os.path.join(self.output_path, f"{os.path.splitext(file_name)[0]}.pdf")
                pdfkit.from_file(html_file_path, output_pdf_path, configuration=config)
                logging.info(f"Converted {file_name} to PDF successfully.")
            
            if os.path.exists(html_file_path):
                os.remove(html_file_path)

        except Exception as e:
            logging.error(f"Error converting Excel to PDF: {e}")
            
    def pdf2excel(self):
        pdf_files = [i for i in os.listdir(self.input_path_dir) if i.lower().endswith(".pdf")]
        if not pdf_files:
            logging.info("No PDF files found in the input directory.")
            return

        for pdf_file in pdf_files:
            pdf_path = os.path.join(self.input_path_dir, pdf_file)
            output_excel = os.path.join(self.output_path, f"{os.path.splitext(pdf_file)[0]}.xlsx")
            try:
                with pdfplumber.open(pdf_path) as pdf:
                    all_tables = []
                    for page in pdf.pages:
                        tables = page.extract_tables()
                        for table in tables:
                            df = pd.DataFrame(table)
                            all_tables.append(df)
                    
                    if all_tables:
                        combined_df = pd.concat(all_tables, ignore_index=True)
                        with pd.ExcelWriter(output_excel) as writer:
                            combined_df.to_excel(writer, index=False, sheet_name='Sheet1')
                        logging.info(f"Converted {pdf_file} to Excel successfully.")
                    else:
                        logging.info(f"No tables found in {pdf_file}.")
            except Exception as e:
                logging.error(f"Error during conversion of {pdf_file}: {e}")
            
    def pdf_to_pptx(self):
        try:
            if not os.path.exists(self.input_path_dir):
                return
            
            input_files = [f for f in os.listdir(self.input_path_dir) if f.lower().endswith(".pdf")]
            if not input_files:
                return

            for pdf_file in input_files:
                pdf_path = os.path.join(self.input_path_dir, pdf_file)  # Full path to the PDF file
                pptx_file = os.path.join(self.output_path, os.path.splitext(os.path.basename(pdf_file))[0] + ".pptx")
                presentation = Presentation()
                document = None
                image_paths = []

                try:
                    document = fitz.open(pdf_path)

                    for page_number in range(len(document)):
                        page = document.load_page(page_number)
                        pixmap = page.get_pixmap()

                        image_path = os.path.join(self.output_path, f'page_{page_number + 1}.png')
                        image_paths.append(image_path)

                        pixmap.save(image_path)

                        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height)

                    presentation.save(pptx_file)

                except Exception as e:
                    logging.error(f"Error processing PDF file {pdf_file}: {e}")

                finally:
                    if document:
                        document.close()

                    for image_path in image_paths:
                        if os.path.exists(image_path):
                            try:
                                os.remove(image_path)
                            except Exception as e:
                                logging.error(f"Error removing temporary image {image_path}: {e}")

        except Exception as e:
            logging.error(f"Error during PDF to PPTX conversion: {e}")

    def pptx2pdf(self):
        try:
            pptx_files = [i for i in os.listdir(self.input_path_dir) if i.lower().endswith(".pptx")]
            
            if not pptx_files:
                logging.info("No PPTX files found in the input directory.")
                return

            for pptx_file in pptx_files:
                input_pptx = os.path.join(self.input_path_dir, pptx_file)
                output_pdf = os.path.join(self.output_path, f"{pptx_file[:-5]}.pdf")

                powerpoint = client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = 1  # Visible = 1 means the application window is visible 0 error

                try:
                    presentation = powerpoint.Presentations.Open(os.path.abspath(input_pptx))
                    
                    presentation.SaveAs(os.path.abspath(output_pdf), 32)  # 32 is the format type for PDF

                    presentation.Close()
                except pywintypes.com_error as e:
                    logging.error(f"Error opening or saving presentation '{pptx_file}': {e}")
                finally:
                    powerpoint.Quit()
        
        except Exception as e:
            logging.error(f"Error during PPTX to PDF conversion: {e}")

    def image2pdf(self):
        img_files = [os.path.join(self.input_path_dir, i) for i in os.listdir(self.input_path_dir) if i.lower().endswith(".jpg") or i.lower().endswith(".png")]

        if not img_files:
            logging.info("No images found to convert.")
            return
        output_pdf = os.path.join(self.output_path, f'{os.path.splitext(os.path.basename(img_files[0]))[0]}.pdf')

        try:
            with open(output_pdf, "wb") as f:
                f.write(img2pdf.convert(img_files))

            logging.info(f"Images successfully converted to {output_pdf}")
        except Exception as e:
            logging.error(f"Error during image to PDF conversion: {e}")
  
    def pdf2images(self):
        try:
            pdf_files = [i for i in os.listdir(self.input_path_dir) if i.lower().endswith(".pdf")]
            if not pdf_files:
                logging.info("No PDF files found in the input directory.")
                return
            
            for pdf_file in pdf_files:
                pdf_path = os.path.join(self.input_path_dir, pdf_file)
                if not os.path.exists(pdf_path):
                    logging.info(f"File not found: {pdf_path}")
                    continue

                try:
                    document = fitz.open(pdf_path)
                    for page_number in range(len(document)):
                        page = document.load_page(page_number)
                        pixmap = page.get_pixmap()
                        image_path = os.path.join(self.output_path, f"page_{page_number + 1}.png")
                        pixmap.save(image_path)
                        logging.info(f"Saved page {page_number + 1} as {image_path}")
                    document.close()
                except Exception as e:
                    logging.error(f"Error processing {pdf_file}: {e}")

                finally:
                    if 'document' in locals() and document:
                        document.close()
        except Exception as e:
            logging.error(f"Error during PDF to image conversion: {e}")
